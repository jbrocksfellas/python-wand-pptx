# Install (pptx, wand) python modules and ImageMagick software to use this script.
from os import listdir
from io import BytesIO
from wand.image import Image
from wand.drawing import Drawing
from pptx import Presentation
from pptx.util import Inches

# creating Image Object for logo
logo = Image(filename="./nike_black.png")
logoRatio = logo.width/logo.height
logoHeight = 100
logo.sample(int(logoRatio * logoHeight), logoHeight)
# path to access the images for presentation
path = "./testing images/"
# path to save presentation (pptx)
savingPath = "./testing images/"
# list of filename in the path
dirList = listdir(path)
# supported format list, you can add more formats.
supported_formats = ['.jpg']
prs = Presentation()

# createPPT funtion creates the ppt
def createPPT(binImg, i):
	# insert binary image into the memory
	f = BytesIO(binImg)
	# Add a slide to the Presentation
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	# Add title
	slide.shapes.title.text = f"Heading {i}"
	# Add subtitle
	slide.placeholders[1].text = f"Subheading {i}"
	# Image size and position using Inches class
	left = Inches(1)
	top = Inches(2.5)
	height = Inches(4)
	# Add picture to the slide
	slide.shapes.add_picture(f, left, top, height=height)

def insertLogo():
	# Tracking variable for the slide number
	i = 1
	# loop for multiple images
	for fname in dirList:
		# loop for distinct formats
		for supFormat in supported_formats:
			# condition if file has the extension (jpeg)
			if supFormat in fname:
				# Creating Image object for current filename
				with Image(filename=path+fname) as img:
					# Keeping aspect ratio
					ratio = img.width/img.height
					# set desired height
					height = 1000
					# setting (width, height) of image w.r.t above height and ratio
					img.sample(int(ratio * height), height)
					with Drawing() as draw:
						draw.composite(operator="src_atop", left=20, top=20, width=logo.width, height=logo.height, image=logo)
						draw(img)
						img.format = 'jpeg'
						# create binary code for image
						jpeg_bin = img.make_blob()
						createPPT(jpeg_bin, i)
						i+=1
	prs.save(f"{savingPath}Indycium.pptx")

insertLogo()


